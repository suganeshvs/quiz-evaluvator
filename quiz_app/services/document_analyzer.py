import os
import re
from pypdf import PdfReader
from pptx import Presentation

class DocumentAnalyzer:
    """
    Service responsible for parsing PDF and PPT/PPTX files,
    extracting text, detecting diagrams/images, splitting content into
    structured paragraphs and lines, and classifying content quantity.
    """

    @staticmethod
    def process_document(document_obj):
        """
        Processes a Document model instance and creates DocumentPage instances for it.
        """
        document_obj.processing_status = 'PROCESSING'
        document_obj.save()

        try:
            file_path = document_obj.file.path if document_obj.file else None
            pages_data = []

            if file_path and os.path.exists(file_path):
                ext = os.path.splitext(file_path)[1].lower()
                if ext == '.pdf':
                    pages_data = DocumentAnalyzer._parse_pdf(file_path)
                elif ext in ['.ppt', '.pptx']:
                    pages_data = DocumentAnalyzer._parse_pptx(file_path)

            # Fallback if no file uploaded or text empty (e.g. for synthetic/demo doc)
            if not pages_data:
                pages_data = DocumentAnalyzer._generate_sample_pages(document_obj.title)

            # Save pages to database
            document_obj.pages.all().delete()  # clear old pages if re-processing
            for page_info in pages_data:
                document_obj.pages.create(
                    page_number=page_info['page_number'],
                    extracted_text=page_info['extracted_text'],
                    content_quantity=page_info['content_quantity'],
                    has_image=page_info['has_image'],
                    image_description=page_info['image_description'],
                    important_topics=page_info['important_topics'],
                    paragraphs_data=page_info['paragraphs_data']
                )

            document_obj.total_pages = len(pages_data)
            document_obj.processing_status = 'COMPLETED'
            document_obj.save()
            return True

        except Exception as e:
            print(f"Error processing document {document_obj.id}: {e}")
            document_obj.processing_status = 'FAILED'
            document_obj.save()
            return False

    @staticmethod
    def _parse_pdf(file_path):
        reader = PdfReader(file_path)
        pages_data = []

        for idx, page in enumerate(reader.pages):
            page_num = idx + 1
            text = page.extract_text() or ""

            # Detect images/images count in pypdf
            has_image = len(page.images) > 0 if hasattr(page, 'images') else False
            image_desc = ""
            if has_image:
                image_desc = f"Diagram / Figure detected on Page {page_num} explaining key visual concept."
            elif re.search(r'(figure|diagram|image|illustration)\s+\d+', text, re.IGNORECASE):
                has_image = True
                image_desc = f"Educational diagram referenced on Page {page_num}."

            paragraphs_data, content_qty, topics = DocumentAnalyzer._structure_text(text, page_num, has_image)

            pages_data.append({
                'page_number': page_num,
                'extracted_text': text,
                'content_quantity': content_qty,
                'has_image': has_image,
                'image_description': image_desc,
                'important_topics': topics,
                'paragraphs_data': paragraphs_data
            })

        return pages_data

    @staticmethod
    def _parse_pptx(file_path):
        prs = Presentation(file_path)
        pages_data = []

        for idx, slide in enumerate(prs.slides):
            page_num = idx + 1
            text_runs = []
            has_image = False

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_runs.append(shape.text_frame.text)
                if shape.shape_type == 13 or hasattr(shape, 'image'):  # MSO_SHAPE_TYPE.PICTURE
                    has_image = True

            full_text = "\n".join(text_runs)
            image_desc = f"Slide diagram on Slide {page_num} explaining concepts." if has_image else ""

            paragraphs_data, content_qty, topics = DocumentAnalyzer._structure_text(full_text, page_num, has_image)

            pages_data.append({
                'page_number': page_num,
                'extracted_text': full_text,
                'content_quantity': content_qty,
                'has_image': has_image,
                'image_description': image_desc,
                'important_topics': topics,
                'paragraphs_data': paragraphs_data
            })

        return pages_data

    @staticmethod
    def _structure_text(text, page_num, has_image):
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        if not paragraphs:
            paragraphs = [p.strip() for p in text.split('\n') if p.strip()]

        if not paragraphs:
            paragraphs = [f"Page {page_num} educational content and visual figures."]

        paragraphs_data = []
        for p_idx, p_text in enumerate(paragraphs, 1):
            lines = [l.strip() for l in p_text.replace('.', '.\n').split('\n') if l.strip()]
            paragraphs_data.append({
                "paragraph_number": p_idx,
                "text": p_text,
                "lines": lines
            })

        total_words = len(text.split())
        if total_words > 80:
            content_qty = 'HIGH'
        elif total_words > 25:
            content_qty = 'MEDIUM'
        else:
            content_qty = 'LOW'

        # Extract basic topic keywords
        words = set(re.findall(r'\b[A-Za-z]{5,}\b', text))
        topics = list(words)[:6]

        return paragraphs_data, content_qty, topics

    @staticmethod
    def _generate_sample_pages(title="Light Chapter 1"):
        """
        Generates comprehensive 20-page sample educational content on Physics / Light.
        Used for initial demo seeding or when uploaded file is empty.
        """
        sample_topics = [
            # Page 1
            ("Introduction to Light & Energy", 
             "Light is a form of electromagnetic radiation and energy that enables us to perceive the world around us. It travels in straight lines at a speed of 3 x 10^8 meters per second in a vacuum.\n\nLight sources are divided into luminous objects, such as the Sun, light bulbs, and stars, which produce their own light, and non-luminous objects, which reflect light from other sources.\n\nWithout light, optical perception and vision would be impossible.",
             True, "Diagram showing the Sun emitting light rays toward the Earth and non-luminous objects."),
            
            # Page 2
            ("Sources of Light & Rectilinear Propagation",
             "Sources of light can be classified into natural sources like stars and lightning, and artificial sources such as candles and LEDs.\n\nRectilinear propagation of light is the fundamental principle stating that light travels in straight lines in a homogeneous medium.\n\nThis principle explains shadow formation, pinhole camera operations, and solar or lunar eclipses.",
             True, "Diagram showing rectilinear propagation of light passing through three aligned pinhole screens."),
            
            # Page 3
            ("Reflection of Light & Laws of Reflection",
             "Reflection occurs when a ray of light falls on a polished smooth surface and bounces back into the same medium.\n\nThe Laws of Reflection state: 1. The angle of incidence is equal to the angle of reflection (i = r). 2. The incident ray, the reflected ray, and the normal at the point of incidence all lie in the same plane.\n\nSpecular reflection happens on smooth surfaces like mirrors, producing clear images, whereas diffuse reflection occurs on rough surfaces.",
             True, "Ray diagram illustrating incident ray, reflected ray, normal vector, angle of incidence (i), and angle of reflection (r)."),
            
            # Page 4
            ("Types of Mirrors: Plane and Spherical Mirrors",
             "A plane mirror has a flat reflecting surface. Images formed by plane mirrors are virtual, erect, laterally inverted, and equal in size to the object.\n\nSpherical mirrors are part of a reflecting sphere and are divided into concave mirrors and convex mirrors.\n\nConcave mirrors curve inward like a spoon's bowl and can form both real and virtual images depending on object position.",
             False, ""),

            # Page 5
            ("Spherical Mirror Terminology & Focal Length",
             "The Pole (P) is the geometric center of the reflecting surface of a spherical mirror.\n\nThe Center of Curvature (C) is the center of the sphere of which the mirror forms a part. The Principal Focus (F) is the point on the principal axis where light rays parallel to the axis converge or appear to diverge.\n\nThe Focal Length (f) is the distance from the Pole to the Principal Focus. The radius of curvature (R) relates to focal length by R = 2f.",
             True, "Detailed spherical mirror diagram highlighting Pole (P), Focus (F), Center of Curvature (C), and Focal Length (f)."),

            # Page 6
            ("Image Formation by Concave Mirrors",
             "When an object is placed at infinity, a concave mirror forms a real, inverted, highly diminished image at the principal focus.\n\nWhen placed beyond C, the image forms between C and F, real, inverted, and diminished.\n\nWhen placed between the Pole and Focus, the image is virtual, erect, and magnified, which makes concave mirrors suitable for shaving mirrors and dentist instruments.",
             True, "Ray diagram showing concave mirror image formation when object is placed between F and P."),

            # Page 7
            ("Convex Mirrors & Applications",
             "Convex mirrors curve outward toward the light source and always form virtual, erect, and diminished images regardless of object distance.\n\nBecause they provide a wide field of view, convex mirrors are universally used as rear-view mirrors in automobiles and security mirrors in stores.\n\nThey allow drivers to see a much larger traffic area behind them than plane mirrors would allow.",
             True, "Diagram showing convex mirror reflecting diverging rays and forming a diminished virtual image."),

            # Page 8
            ("Mirror Formula & Magnification",
             "The mirror formula gives the mathematical relation between object distance (u), image distance (v), and focal length (f): 1/f = 1/v + 1/u.\n\nLinear Magnification (m) is defined as the ratio of image height (h') to object height (h): m = h'/h = -v/u.\n\nA negative magnification indicates a real image, while a positive magnification indicates a virtual image.",
             False, ""),

            # Page 9
            ("Refraction of Light & Medium Density",
             "Refraction is the bending of light when it passes obliquely from one transparent medium to another of different optical density.\n\nWhen light travels from an optically rarer medium (like air) to a denser medium (like glass), it bends toward the normal line.\n\nWhen light travels from a denser medium to a rarer medium, it bends away from the normal line.",
             True, "Ray diagram showing refraction of light at an air-glass interface bending toward the normal."),

            # Page 10
            ("Snell's Law & Refractive Index",
             "Snell's Law of Refraction states that the ratio of the sine of the angle of incidence to the sine of the angle of refraction is a constant for a given pair of media: sin(i) / sin(r) = n2 / n1.\n\nThe absolute refractive index (n) of a medium is given by n = c / v, where c is the speed of light in vacuum and v is the speed in the medium.\n\nDiamond has a high refractive index of 2.42, causing strong light bending and sparkling brilliance.",
             False, ""),

            # Page 11
            ("Refraction Through a Glass Slab",
             "When light passes through a rectangular glass slab, it undergoes refraction twice: at the air-glass surface and the glass-air surface.\n\nThe emergent ray is parallel to the original incident ray, but shifted sideways by a distance called lateral displacement.\n\nLateral displacement increases with slab thickness, refractive index, and angle of incidence.",
             True, "Diagram of rectangular glass slab showing incident ray, refracted ray, emergent ray, and lateral displacement."),

            # Page 12
            ("Total Internal Reflection & Critical Angle",
             "Total Internal Reflection (TIR) occurs when light traveling in a denser medium strikes a rarer medium boundary at an angle greater than the critical angle.\n\nThe Critical Angle is the angle of incidence in a denser medium for which the angle of refraction in the rarer medium is exactly 90 degrees.\n\nTIR is used in optical fibers for internet communications, endoscopes in medicine, and causes mirages in hot deserts.",
             True, "Diagram showing light rays at increasing angles leading to critical angle refraction and total internal reflection."),

            # Page 13
            ("Lenses: Convex and Concave Lenses",
             "A lens is a transparent optical medium bounded by two spherical surfaces.\n\nA convex lens is thicker at the center than at the edges and acts as a converging lens for light rays.\n\nA concave lens is thinner at the center than at the edges and acts as a diverging lens for light rays.",
             True, "Diagram comparing light rays converging through a convex lens versus diverging through a concave lens."),

            # Page 14
            ("Lens Formula & Power of Lens",
             "The Lens Formula states: 1/f = 1/v - 1/u, connecting object distance (u), image distance (v), and focal length (f).\n\nPower of a Lens (P) is defined as the reciprocal of its focal length in meters: P = 1 / f(in meters).\n\nThe SI unit of lens power is the Dioptre (D). Convex lenses have positive power while concave lenses have negative power.",
             False, ""),

            # Page 15
            ("Human Eye Structure & Vision Mechanism",
             "The human eye acts like a camera using a flexible convex lens to focus light onto the light-sensitive retina at the back.\n\nThe cornea is the outer transparent layer that performs most of the light refraction. The iris controls the size of the pupil to regulate light intensity.\n\nThe ciliary muscles change the curvature and focal length of the eye lens to focus on objects at varying distances (Power of Accommodation).",
             True, "Detailed anatomical diagram of the human eye showing cornea, iris, pupil, eye lens, ciliary muscles, and retina."),

            # Page 16
            ("Defects of Vision: Myopia and Hypermetropia",
             "Myopia (short-sightedness) is a condition where a person can see near objects clearly but distant objects appear blurry because images form in front of the retina.\n\nMyopia is corrected using a concave lens of suitable focal length.\n\nHypermetropia (far-sightedness) occurs when images form behind the retina, corrected using a convex lens.",
             True, "Diagram showing myopic eye image forming before retina and its correction using a concave lens."),

            # Page 17
            ("Dispersion of Light Through a Prism",
             "Dispersion is the splitting of white light into its component seven colors (VIBGYOR: Violet, Indigo, Blue, Green, Yellow, Orange, Red) when passing through a glass prism.\n\nDispersion occurs because different colors travel at different speeds in glass, with Violet bending the most and Red bending the least.\n\nRecombination of spectrum colors using a second inverted prism reproduces pure white light.",
             True, "Prism diagram illustrating white light ray splitting into seven spectrum colors (VIBGYOR)."),

            # Page 18
            ("Rainbow Formation in Nature",
             "A rainbow is a natural spectrum appearing in the sky after rain, caused by dispersion, refraction, and internal reflection of sunlight by tiny water droplets.\n\nWater droplets act like small glass prisms. Sunlight refracts upon entering the drop, internally reflects off the back surface, and refracts again upon exiting.\n\nRed light emerges at an angle of 42 degrees while Violet emerges at 40 degrees relative to the observer line.",
             True, "Diagram showing sunlight entering a raindrop, undergoing refraction, internal reflection, and color dispersion."),

            # Page 19
            ("Atmospheric Refraction & Twinkling Stars",
             "Atmospheric refraction is the refraction of light caused by Earth's atmosphere having layers of varying air density and refractive index.\n\nStars twinkle because light from stars passes through continuously fluctuating atmospheric layers, causing the apparent position and brightness to fluctuate.\n\nPlanets do not twinkle because they are much closer to Earth and act as extended light sources.",
             False, ""),

            # Page 20
            ("Scattering of Light & Tyndall Effect",
             "Scattering of light is the phenomenon where fine particles suspend in a medium diffuse light rays in all directions.\n\nThe Tyndall Effect occurs when a light beam passes through a colloidal solution, making the path of light visible.\n\nThe sky appears blue because short wavelength blue light is scattered much more strongly by fine atmospheric molecules than longer red light.",
             True, "Diagram illustrating Tyndall Effect in a colloid container with a visible light beam path.")
        ]

        pages_data = []
        for p_num, (heading, text_body, has_img, img_desc) in enumerate(sample_topics, 1):
            paragraphs = [p.strip() for p in text_body.split('\n\n') if p.strip()]
            paragraphs_data = []
            for p_idx, p_text in enumerate(paragraphs, 1):
                lines = [l.strip() for l in p_text.replace('.', '.\n').split('\n') if l.strip()]
                paragraphs_data.append({
                    "paragraph_number": p_idx,
                    "text": p_text,
                    "lines": lines
                })

            full_text = f"Chapter: {title}\nPage {p_num}: {heading}\n\n" + text_body

            pages_data.append({
                'page_number': p_num,
                'extracted_text': full_text,
                'content_quantity': 'HIGH',
                'has_image': has_img,
                'image_description': img_desc,
                'important_topics': [heading, "Physics", "Light", f"Topic-{p_num}"],
                'paragraphs_data': paragraphs_data
            })

        return pages_data
